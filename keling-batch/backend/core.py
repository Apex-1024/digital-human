"""课灵 AI 批量制课系统 — 后端核心"""
import os
import re
import sys
import json
import time
import uuid
import shutil
import threading
import subprocess
import logging
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import List, Optional, Dict, Any
from concurrent.futures import ThreadPoolExecutor

# MinerU 模型源：用 modelscope 避免代理问题
os.environ.setdefault("MINERU_MODEL_SOURCE", "modelscope")

BASE_DIR = Path(__file__).resolve().parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "output"
ASSETS_DIR = BASE_DIR / "assets"
AVATAR_DIR = ASSETS_DIR / "avatars"
BGM_DIR = ASSETS_DIR / "bgm"
JOBS_FILE = BASE_DIR / "jobs.json"
# 调试目录：每次任务记录 PPT 提取 / 教案提取 / LLM 文案合成结果
DEBUG_DIR = BASE_DIR / "test_debug"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
DEBUG_DIR.mkdir(parents=True, exist_ok=True)
AVATAR_DIR.mkdir(parents=True, exist_ok=True)
BGM_DIR.mkdir(parents=True, exist_ok=True)

logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s")
log = logging.getLogger("keling-batch")


def _find_ffmpeg() -> str:
    try:
        import imageio_ffmpeg
        path = imageio_ffmpeg.get_ffmpeg_exe()
        if Path(path).exists():
            log.info(f"使用 ffmpeg: {path}")
            return path
    except ImportError:
        pass
    return "ffmpeg"


FFMPEG_BIN = _find_ffmpeg()

_p2t_engine = None


def _get_p2t():
    """懒加载 MinerU 引擎（PDF 解析，含公式识别）。

    MinerU 3.4 是异步 API 模式：本地无 api_url 时会自动启动临时 FastAPI 服务。
    这里不预创建服务，只确认 mineru 包可导入。真正调用在 _ocr_pdf_with_mineru。
    """
    global _p2t_engine
    if _p2t_engine is not None:
        return _p2t_engine
    if _p2t_engine is False:
        return None
    try:
        import mineru  # noqa: F401
        _patch_pdfium_image_get_pos()
        _p2t_engine = True
        log.info("MinerU 包可用")
    except Exception as e:
        log.warning(f"MinerU 包不可用：{e}")
        _p2t_engine = False
    return _p2t_engine


def _patch_pdfium_image_get_pos():
    """兼容补丁：pypdfium2 5.x 移除了 PdfImage.get_pos()，MinerU 3.4.1 仍依赖它。

    pdf_classify.get_high_image_coverage_ratio_pdfium 调用 page_object.get_pos()
    获取图像四元组 (left, bottom, right, top)。pypdfium2 5.x 改名为 get_bounds()，
    返回格式完全一致。

    本函数打主进程补丁；MinerU 子进程补丁通过 PYTHONPATH + usercustomize.py 注入
    （见 _ensure_mineru_subprocess_patch）。
    """
    try:
        import pypdfium2 as _pdfium
        from pypdfium2._helpers.pageobjects import PdfImage
        if not hasattr(PdfImage, "get_pos") and hasattr(PdfImage, "get_bounds"):
            PdfImage.get_pos = lambda self: self.get_bounds()
            log.info("已为 PdfImage.get_pos 打主进程兼容补丁（委托给 get_bounds）")
    except Exception as e:
        log.warning(f"PdfImage.get_pos 主进程补丁失败（忽略，将由 MinerU except 兜底）：{e}")


# 标记子进程补丁是否已注入环境变量，避免重复设置
_mineru_subprocess_patch_injected = False


def _ensure_mineru_subprocess_patch():
    """注入 PYTHONPATH 让 MinerU LocalAPIServer 子进程启动时执行补丁。

    MinerU 的 LocalAPIServer.start() 用 subprocess.Popen 启动独立 Python 进程
    跑 fast_api 模块，env=os.environ.copy() 继承主进程环境变量。主进程的
    monkey-patch 不会传到子进程，但 PYTHONPATH 会。子进程启动时 site 模块会
    自动 import PYTHONPATH 目录里的 usercustomize.py，从而打上补丁。
    """
    global _mineru_subprocess_patch_injected
    if _mineru_subprocess_patch_injected:
        return
    patch_dir = str(Path(__file__).resolve().parent.parent / "_mineru_patch")
    if not Path(patch_dir).is_dir():
        log.warning(f"MinerU 子进程补丁目录不存在：{patch_dir}")
        return
    existing = os.environ.get("PYTHONPATH", "")
    sep = ";" if os.name == "nt" else ":"
    if patch_dir in existing.split(sep):
        _mineru_subprocess_patch_injected = True
        return
    os.environ["PYTHONPATH"] = (
        f"{patch_dir}{sep}{existing}" if existing else patch_dir
    )
    _mineru_subprocess_patch_injected = True
    log.info(f"已注入 PYTHONPATH 供 MinerU 子进程加载补丁：{patch_dir}")


def _ocr_pdf_with_mineru(pdf_path: str) -> List[List[tuple]]:
    """用 MinerU 解析整本 PDF，按页返回 [(text, kind, y), ...]。

    - text/kind: 'text' 或 'formula'
    - y: block 的 bbox 顶部坐标（用于排序），来自 MinerU JSON 的 bbox[1]
    - 输出按 page_idx 拆分，result[i] 对应第 i 页
    """
    if not _get_p2t():
        return []

    import asyncio
    import httpx
    import tempfile
    from mineru.cli import api_client as _api_client

    # 确保 MinerU 子进程启动时加载 PdfImage.get_pos 补丁
    _ensure_mineru_subprocess_patch()

    pdf_path = Path(pdf_path)
    output_dir = pdf_path.parent / f"_mineru_{pdf_path.stem}"
    output_dir.mkdir(exist_ok=True)

    form_data = _api_client.build_parse_request_form_data(
        lang_list=["ch"],
        backend="pipeline",
        parse_method="auto",
        formula_enable=True,
        table_enable=False,
        image_analysis=False,
        server_url=None,
        start_page_id=0,
        end_page_id=None,
        return_md=False,            # 不需要 md，只要中间 json
        return_middle_json=True,
        return_model_output=False,
        return_content_list=False,
        return_images=False,
        response_format_zip=True,
        return_original_file=False,
    )
    upload = [_api_client.UploadAsset(path=str(pdf_path), upload_name=pdf_path.name)]

    async def _run():
        async with httpx.AsyncClient(
            timeout=_api_client.build_http_timeout(),
            follow_redirects=True,
        ) as http:
            local_server = _api_client.LocalAPIServer()
            base_url = local_server.start()
            try:
                health = await _api_client.wait_for_local_api_ready(http, local_server)
                submit = await _api_client.submit_parse_task(
                    base_url=health.base_url,
                    upload_assets=upload,
                    form_data=form_data,
                )
                await _api_client.wait_for_task_result(
                    client=http,
                    submit_response=submit,
                    task_label=pdf_path.name,
                    status_snapshot_callback=lambda s: None,
                )
                zip_path = await _api_client.download_result_zip(
                    client=http,
                    submit_response=submit,
                    task_label=pdf_path.name,
                )
                try:
                    _api_client.safe_extract_zip(zip_path, output_dir)
                finally:
                    zip_path.unlink(missing_ok=True)
            finally:
                local_server.stop()

    try:
        asyncio.run(_run())
    except Exception as e:
        log.warning(f"MinerU 解析失败 {pdf_path}：{e}")
        return []

    # 找到 middle_json
    json_path = None
    for cand in output_dir.rglob("*_middle.json"):
        json_path = cand
        break
    if not json_path:
        log.warning(f"MinerU 未输出 middle.json：{output_dir}")
        return []

    import json
    try:
        data = json.loads(json_path.read_text(encoding="utf-8"))
    except Exception as e:
        log.warning(f"MinerU JSON 解析失败：{e}")
        return []

    pdf_info = data.get("pdf_info", [])
    pages: List[List[tuple]] = []

    def _extract_spans(blocks: List[dict]) -> List[tuple]:
        items: List[tuple] = []
        for b in blocks:
            btype = b.get("type", "")
            # 跳过图片、表格等非文本块
            if btype in ("image", "image_body", "table", "table_body"):
                continue
            bbox = b.get("bbox", [0, 0, 0, 0])
            y_top = bbox[1] if len(bbox) >= 4 else 0
            lines = b.get("lines", [])
            if not lines:
                # 有些 block 直接有 content
                content = b.get("content", "").strip()
                if content:
                    items.append((content, "text", y_top))
                continue
            for line in lines:
                for span in line.get("spans", []):
                    stype = span.get("type", "text")
                    content = span.get("content", "").strip()
                    if not content:
                        continue
                    if stype == "inline_equation":
                        items.append((content, "formula", y_top))
                    elif stype == "text":
                        items.append((content, "text", y_top))
                    # 其他类型（如 display_formula）也按公式处理
                    elif "equation" in stype or "formula" in stype:
                        items.append((content, "formula", y_top))
        # 按 y 排序
        items.sort(key=lambda x: x[2])
        return items

    # 按 page_idx 排序，保证页序正确
    pdf_info_sorted = sorted(pdf_info, key=lambda p: p.get("page_idx", 0))
    for page in pdf_info_sorted:
        blocks = page.get("para_blocks") or page.get("preproc_blocks") or []
        pages.append(_extract_spans(blocks))

    return pages


# ============ LLM 口播转换（DeepSeek）============
_LLM_CACHE: Dict[str, str] = {}


def latex_to_chinese_llm(latex: str) -> str:
    """用 DeepSeek 把 LaTeX 公式转成中文口播读法，失败回退到正则方案"""
    if not latex:
        return ""
    s = latex.strip().strip("$").strip()
    if not s:
        return ""

    # 命中缓存直接返回
    if s in _LLM_CACHE:
        return _LLM_CACHE[s]

    api_key = os.environ.get("DEEPSEEK_API_KEY", "").strip()
    if not api_key:
        # 未配置 key，直接走正则回退
        return latex_to_chinese(latex)

    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com/v1", timeout=15)
        resp = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {
                    "role": "system",
                    "content": (
                        "你是数学公式口播转换器。输入 LaTeX 公式，输出自然的中文口播读法。"
                        "规则：1) 只输出中文口播文本，不加任何解释或前后缀；"
                        "2) 积分 ∫_a^b 读作'从 a 到 b 的积分'；"
                        "3) 求和 Σ_{i=1}^{n} 读作'i 从 1 到 n 求和'；"
                        "4) 分数 \\frac{a}{b} 读作'b 分之 a'；"
                        "5) 根号 \\sqrt{x} 读作'根号下 x'；"
                        "6) 上标 a^{2} 读作'a 的 2 次方'；"
                        "7) 单字符下标 T_{1} 直接读作'T1'（拼接，不读'下标'）；"
                        "8) 多字符下标 T_{ij} 读作'T 下标 ij'；"
                        "9) 区间 [a, b] 读作'区间 a 到 b'；"
                        "10) 希腊字母用中文名（α→阿尔法，β→贝塔等）；"
                        "11) \\geq 读'大于等于'，\\leq 读'小于等于'，\\neq 读'不等于'；"
                        "12) \\cdot 读'乘'，\\times 读'乘以'，\\div 读'除以'；"
                        "13) \\to 读'趋近于'，\\rightarrow 读'趋向'；"
                        "14) \\infty 读'无穷'，\\partial 读'偏导'；"
                        "15) 保留变量字母原文（如 v(t) 读'v t'，不翻译成中文）；"
                        "16) 去掉所有 LaTeX 命令（\\left \\right \\! \\, 等），只保留口播内容。"
                    ),
                },
                {"role": "user", "content": s},
            ],
            temperature=0.0,
            max_tokens=256,
        )
        result = resp.choices[0].message.content.strip()
        if result:
            _LLM_CACHE[s] = result
            return result
    except Exception as e:
        log.warning(f"DeepSeek LLM 转换失败，回退正则：{e}")

    # 回退到正则方案
    return latex_to_chinese(latex)


def latex_to_chinese(latex: str) -> str:
    """LaTeX 公式 → 中文口播读法"""
    if not latex:
        return ""
    s = latex.strip().strip("$").strip()

    # 希腊字母
    # 规范化：MinerU 输出 LaTeX 时在符号间插入大量空格（如 "T _ { 1 }"）
    # 去除 \ _ ^ { } , 周围的空格，使正则能稳定匹配
    s = re.sub(r"\s*([\\_^{},])\s*", r"\1", s)

    greek = {
        r"\alpha": "阿尔法", r"\beta": "贝塔", r"\gamma": "伽马",
        r"\delta": "德尔塔", r"\epsilon": "艾普西隆", r"\varepsilon": "艾普西隆",
        r"\zeta": "泽塔", r"\eta": "伊塔", r"\theta": "西塔", r"\vartheta": "西塔",
        r"\iota": "约塔", r"\kappa": "卡帕", r"\lambda": "拉姆达",
        r"\mu": "缪", r"\nu": "纽", r"\xi": "克西",
        r"\pi": "派", r"\varpi": "派", r"\rho": "柔", r"\varrho": "柔",
        r"\sigma": "西格玛", r"\varsigma": "西格玛",
        r"\tau": "陶", r"\upsilon": "宇普西隆", r"\phi": "弗爱", r"\varphi": "弗爱",
        r"\chi": "凯", r"\psi": "普西", r"\omega": "欧米伽",
        r"\Gamma": "伽马", r"\Delta": "德尔塔", r"\Theta": "西塔",
        r"\Lambda": "拉姆达", r"\Xi": "克西", r"\Pi": "派",
        r"\Sigma": "西格玛", r"\Phi": "弗爱", r"\Psi": "普西", r"\Omega": "欧米伽",
    }
    for k, v in greek.items():
        s = s.replace(k, v)

    # 极限：\lim_{x \to a} → x 趋近于 a 时的极限
    s = re.sub(r"\\lim_\{?([^}\\]+)\s*\\to\s*([^}\\]+)\}?",
               r"极限 \1 趋近于 \2 ", s)
    s = s.replace(r"\lim", "极限 ")
    s = s.replace(r"\to", " 趋近于 ")

    # 预处理：单字符下标拍平（如 T_{1} → T1），消除嵌套大括号
    # 这样积分/求和的正则能正确匹配 _{T1}^{T2}（否则 _{T_{1}} 嵌套匹配失败）
    s = re.sub(r"_\{([0-9A-Za-z])\}", r"\1", s)

    # 积分：\int_{a}^{b} → 从 a 到 b 的积分（预处理后嵌套大括号已拍平）
    m = re.search(r"\\int(?:_\{([^{}]+)\})?(?:\^\{([^{}]+)\})?", s)
    if m:
        lo, hi = m.group(1), m.group(2)
        repl = "积分"
        if lo and hi:
            repl = f"从 {lo.strip()} 到 {hi.strip()} 的积分"
        elif lo:
            repl = f"下限 {lo.strip()} 的积分"
        elif hi:
            repl = f"上限 {hi.strip()} 的积分"
        s = s.replace(m.group(0), repl, 1)
    s = s.replace(r"\int", "积分")
    s = s.replace(r"\,dx", " 微元 d x").replace(r"\,dt", " 微元 d t")
    s = s.replace(r"\,dy", " 微元 d y").replace(" dx", " 微元 d x").replace(" dt", " 微元 d t")

    # 求和：\sum_{i=1}^{n} → i 从 1 到 n 求和
    m = re.search(r"\\sum_\{?([^{_}^]+)\}?(?:\^\{?([^{_}^]+)\}?)?", s)
    if m:
        lo, hi = m.group(1), m.group(2)
        repl = "求和"
        if lo and hi:
            repl = f"{lo.strip()} 到 {hi.strip()} 求和"
        s = s.replace(m.group(0), repl, 1)
    s = s.replace(r"\sum", "求和")

    # 乘积：\prod
    s = s.replace(r"\prod", "求积")

    # 分数：\frac{a}{b} → b 分之 a
    def _frac(m):
        a, b = m.group(1), m.group(2)
        return f" {b} 分之 {a} "
    s = re.sub(r"\\frac\{([^{}]*)\}\{([^{}]*)\}", _frac, s)
    s = re.sub(r"\\dfrac\{([^{}]*)\}\{([^{}]*)\}", _frac, s)
    s = re.sub(r"\\tfrac\{([^{}]*)\}\{([^{}]*)\}", _frac, s)

    # 二项式系数
    def _binom(m):
        a, b = m.group(1), m.group(2)
        return f" 从 {b} 中取 {a} "
    s = re.sub(r"\\binom\{([^{}]*)\}\{([^{}]*)\}", _binom, s)

    # 根号：\sqrt{x} / \sqrt[n]{x}
    m = re.search(r"\\sqrt\[(.*?)\]\{(.*?)\}", s)
    if m:
        s = s.replace(m.group(0), f" {m.group(1)} 次根号下 {m.group(2)} ", 1)
    s = re.sub(r"\\sqrt\{([^{}]*)\}", r" 根号下 \1 ", s)
    s = s.replace(r"\sqrt", "根号")

    # 上下标：a^{bc} → a 的 bc 次方；a_{ij} → a 下标 ij
    # 单字符下标（如 T_{2} → T2）直接拼接，更自然；多字符下标才读"下标"
    s = re.sub(r"\^\{([^{}]+)\}", r" 的 \1 次方", s)
    s = re.sub(r"_\{([0-9A-Za-z])\}", r"\1", s)  # 单字符下标直接拼接
    s = re.sub(r"_\{([^{}]+)\}", r" 下标 \1", s)  # 多字符下标读"下标"
    s = re.sub(r"\^([0-9A-Za-z])", r" 的 \1 次方", s)
    s = re.sub(r"_([0-9A-Za-z])", r"\1", s)  # 单字符下标直接拼接

    # 关系符
    rel = {
        r"\leq": " 小于等于 ", r"\le": " 小于等于 ",
        r"\geq": " 大于等于 ", r"\ge": " 大于等于 ",
        r"\neq": " 不等于 ", r"\ne": " 不等于 ",
        r"\approx": " 约等于 ", r"\equiv": " 恒等于 ",
        r"\in": " 属于 ", r"\notin": " 不属于 ",
        r"\subset": " 包含于 ", r"\subseteq": " 包含于 ",
        r"\cup": " 并集 ", r"\cap": " 交集 ",
        r"\forall": " 任意 ", r"\exists": " 存在 ",
        r"\infty": " 无穷 ", r"\partial": " 偏导 ",
        r"\nabla": " 梯度 ", r"\pm": " 正负 ", r"\mp": " 负正 ",
        r"\times": " 乘以 ", r"\div": " 除以 ", r"\cdot": " 乘 ",
        r"\rightarrow": " 趋向 ", r"\to": " 趋向 ",
        r"\Rightarrow": " 推出 ", r"\Leftrightarrow": " 等价于 ",
        r"\langle": " 左尖 ", r"\rangle": " 右尖 ",
        r"\|": " 平行 ", r"\perp": " 垂直 ",
    }
    for k, v in rel.items():
        # 加 (?![a-zA-Z]) 防止 \le 匹配到 \left、\in 匹配到 \infty 等
        s = re.sub(re.escape(k) + r"(?![a-zA-Z])", v, s)

    # 算符
    s = s.replace(r"\log", " 对数 ").replace(r"\ln", " 自然对数 ")
    s = s.replace(r"\sin", " 正弦 ").replace(r"\cos", " 余弦 ")
    s = s.replace(r"\tan", " 正切 ").replace(r"\cot", " 余切 ")
    s = s.replace(r"\exp", " 指数 ").replace(r"\max", " 最大值 ")
    s = s.replace(r"\min", " 最小值 ").replace(r"\det", " 行列式 ")

    # 字体修饰去掉（取内容），两轮以支持 \overline{{S}} 这种嵌套大括号
    for cmd in (r"\mathrm", r"\mathbb", r"\mathbf", r"\mathit",
                r"\boldsymbol", r"\overline", r"\bar", r"\vec",
                r"\text", r"\mathcal", r"\mathsf"):
        s = re.sub(re.escape(cmd) + r"\{([^{}]*)\}", r"\1", s)
    # 拍平单层大括号（消除 \overline{{S}} 的内层大括号）
    s = re.sub(r"\{([^{}]+)\}", r"\1", s)
    # 第二轮字体修饰（处理拍平后暴露的命令）
    for cmd in (r"\mathrm", r"\mathbb", r"\mathbf", r"\mathit",
                r"\boldsymbol", r"\overline", r"\bar", r"\vec",
                r"\text", r"\mathcal", r"\mathsf"):
        s = re.sub(re.escape(cmd) + r"\{([^{}]*)\}", r"\1", s)

    # 区间：\left[ a , b \right] → 区间 a 到 b
    m = re.search(r"\\left\[\s*([^,]+?)\s*,\s*([^,]+?)\s*\\right\]", s)
    if m:
        s = s.replace(m.group(0), f" 区间 {m.group(1).strip()} 到 {m.group(2).strip()} ", 1)

    # 括号：\left( \right) → ( )；其他 \left \right 配对简化
    s = s.replace(r"\left(", "(").replace(r"\right)", ")")
    s = s.replace(r"\left[", "[").replace(r"\right]", "]")
    s = s.replace(r"\left|", "|").replace(r"\right|", "|")
    s = s.replace(r"\left\{", "{").replace(r"\right\}", "}")
    s = s.replace(r"\left", " ").replace(r"\right", " ")
    s = s.replace(r"\big", " ").replace(r"\Big", " ")

    # 间距符
    s = s.replace(r"\!", "").replace(r"\,", " ").replace(r"\;", " ").replace(r"\:", " ")
    # \stackrel{a}{b} → b 上方 a
    s = re.sub(r"\\stackrel\{([^{}]*)\}\{([^{}]*)\}", r" \2 上方 \1 ", s)
    s = s.replace(r"\quad", " ").replace(r"\qquad", " ")

    # 残留反斜杠命令直接丢掉
    s = re.sub(r"\\[a-zA-Z]+", " ", s)

    # 多余符号
    s = s.replace("&", " 且 ").replace("\\", " ")
    s = s.replace("^", " 的次方 ").replace("_", " 下标 ")
    s = s.replace("{", " ").replace("}", " ")
    s = s.replace("$", " ")

    # 折叠多空格
    s = re.sub(r"\s+", " ", s).strip()
    return s


def math_symbols_to_chinese(text: str) -> str:
    """数学符号 → 中文口播读法"""
    if not text:
        return ""
    s = text

    # MathType Symbol 字体 Private Use Area 字符映射
    mathtype_map = {
        "\uf021": "±",   # plus-minus
        "\uf022": "≥",   # greaterequal
        "\uf023": "≤",   # lessequal
        "\uf024": "∫",   # integral
        "\uf029": "∞",   # infinity
        "\uf02a": "∂",   # partial
        "\uf02b": "∏",   # product
        "\uf02d": "-",   # minus (Symbol 字体 0x2D)
        "\uf02f": "·",   # middle dot
        "\uf030": "√",   # sqrt
        "\uf034": "∑",   # summation
        "\uf035": "∏",   # product
        "\uf03d": "=",   # equal (Symbol 字体 0x3D)
        "\uf03f": "≈",   # approx
        "\uf040": "≡",   # equiv
        "\uf041": "α",   # alpha
        "\uf042": "β",   # beta
        "\uf043": "×",   # multiply (Symbol 0xB4 → × 实际 0xD7 但映射到这里)
        "\uf045": "ε",   # epsilon
        "\uf046": "φ",   # phi
        "\uf047": "γ",   # gamma
        "\uf048": "η",   # eta
        "\uf049": "ι",   # iota
        "\uf04a": "φ",   # varphi
        "\uf04b": "κ",   # kappa
        "\uf04c": "λ",   # lambda
        "\uf04d": "μ",   # mu
        "\uf04e": "ν",   # nu
        "\uf04f": "ο",   # omicron
        "\uf050": "π",   # pi
        "\uf051": "θ",   # theta
        "\uf052": "ρ",   # rho
        "\uf053": "σ",   # sigma
        "\uf054": "τ",   # tau
        "\uf055": "υ",   # upsilon
        "\uf056": "→",   # arrow
        "\uf057": "ω",   # omega
        "\uf058": "ξ",   # xi
        "\uf059": "ψ",   # psi
        "\uf05a": "ζ",   # zeta
        "\uf0a3": "≤",   # lessequal (Symbol 0xA3)
        "\uf0b3": "≥",   # greaterequal (Symbol 0xB3)
        "\uf0b4": "≠",   # notequal (Symbol 0xB4)
        "\uf0b7": "·",   # centerdot
        "\uf0c5": "×",   # multiply (Symbol 0xC5)
        "\uf0c6": "÷",   # divide (Symbol 0xC6)
        "\uf0d0": "→",   # rightarrow (Symbol 0xD0)
        "\uf0de": "→",   # arrow
        "\uf0e5": "π",   # pi (Symbol 0xE5)
        "\uf0e9": "∇",   # nabla
        "\uf0f4": "√",   # sqrt (Symbol 0xF4)
        "\uf0f5": "∝",   # proportional
        "\uf0f6": "∞",   # infinity
    }
    for k, v in mathtype_map.items():
        s = s.replace(k, v)

    # ---- 复合关系符（先于单字符处理，避免误替换）----
    s = s.replace("≥", " 大于等于 ").replace("≦", " 小于等于 ")
    s = s.replace("≤", " 小于等于 ").replace("≧", " 大于等于 ")
    s = s.replace("≠", " 不等于 ").replace("≈", " 约等于 ")
    s = s.replace("≡", " 恒等于 ").replace("∝", " 正比于 ")
    s = s.replace("→", " 趋向 ").replace("←", " 反向趋向 ")
    s = s.replace("∞", " 无穷 ")
    s = s.replace("±", " 正负 ").replace("∓", " 负正 ")
    s = s.replace("√", " 根号 ").replace("∑", " 求和 ").replace("∏", " 求积 ")
    s = s.replace("∈", " 属于 ").replace("∉", " 不属于 ")
    s = s.replace("∪", " 并集 ").replace("∩", " 交集 ")
    s = s.replace("∀", " 任意 ").replace("∃", " 存在 ")
    s = s.replace("∂", " 偏导 ").replace("∇", " 梯度 ")

    # ---- 区间括号 [a, b] → a 到 b ----
    # [T1,T2] → T1 到 T2 ； [a, b] → a 到 b
    s = re.sub(r"\[\s*([^,\]\[]+?)\s*,\s*([^,\]\[]+?)\s*\]", r" \1 到 \2 ", s)

    # ---- 积分号（如果 OCR 直接识别到 ∫）----
    # ∫_a^b → 从 a 到 b 积分 ； ∫ → 积分
    s = re.sub(r"∫[_\s]*([0-9A-Za-z]+)\s*\^\s*([0-9A-Za-z]+)", r" 从 \1 到 \2 积分 ", s)
    s = s.replace("∫", " 积分 ")

    # ---- 算术运算符 ----
    # 减号：在字母/数字之间用"减"，否则去掉
    s = re.sub(r"([0-9A-Za-z）)])\s*[-−－]\s*([0-9A-Za-z（(])", r"\1 减 \2", s)
    # 加号
    s = re.sub(r"([0-9A-Za-z）)])\s*\+\s*([0-9A-Za-z（(])", r"\1 加 \2", s)
    s = s.replace("×", " 乘以 ").replace("·", " 乘 ").replace("·", " 乘 ")
    s = s.replace("÷", " 除以 ")

    # ---- 等号 ----
    s = s.replace("=", " 等于 ")

    # ---- 大于/小于（在复合符处理后）----
    s = re.sub(r"([0-9A-Za-z）)])\s*>\s*([0-9A-Za-z（(])", r"\1 大于 \2", s)
    s = re.sub(r"([0-9A-Za-z）)])\s*<\s*([0-9A-Za-z（(])", r"\1 小于 \2", s)

    # ---- 上标/下标数字（如 T1, T2, x2）----
    # 单字母+单数字（T1, T2, x2）直接保留，读作"T1""T2"（更自然）
    # 多字符下标才读"下标"（如 T12 → T 下标 12）
    s = re.sub(r"([A-Za-z])([0-9]{2,})", r"\1 下标 \2", s)

    # ---- 撇号（导数 s'(t) → s 撇 (t)）----
    s = s.replace("'", " 撇 ").replace("'", " 撇 ").replace("'", " 撇 ")

    # ---- 百分号 ----
    s = s.replace("%", " 百分之 ")

    # 折叠多空格
    s = re.sub(r"\s+", " ", s).strip()
    # 修正：中文与中文之间的多余空格
    s = re.sub(r"([\u4e00-\u9fff，。：；！？])\s+([\u4e00-\u9fff，。：；！？])", r"\1\2", s)
    return s


# ============ 任务状态 ============
JOB_STATUS_PENDING = "pending"
JOB_STATUS_RUNNING = "running"
JOB_STATUS_DONE = "done"
JOB_STATUS_FAILED = "failed"
JOB_STATUS_CANCELED = "canceled"


def _read_doc_with_com(path: Path) -> str:
    """用 Word COM 读取 .doc 文件文本。需 pywin32 + 已安装 Word，失败抛异常。"""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()
    app = None
    doc = None
    try:
        app = win32com.client.Dispatch("Word.Application")
        app.Visible = False
        doc = app.Documents.Open(str(path.resolve()), ReadOnly=True)
        return doc.Content.Text
    finally:
        try:
            if doc:
                doc.Close(False)
        except Exception:
            pass
        try:
            if app:
                app.Quit()
        except Exception:
            pass


def extract_lesson_plan(path: str) -> str:
    """从教案文件提取纯文本。支持 txt/md/doc/docx，失败返回空串。

    txt/md 直接读文本（utf-8 失败回退 gbk）；docx 用 python-docx 提取段落；
    doc 用 Word COM 读取（需 pywin32 + Word）。
    截断到 8000 字避免 LLM prompt 过长。
    """
    p = Path(path)
    if not p.exists():
        log.warning(f"教案文件不存在：{path}")
        return ""
    ext = p.suffix.lower()
    try:
        if ext in (".txt", ".md"):
            try:
                text = p.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                text = p.read_text(encoding="gbk")
        elif ext == ".docx":
            from docx import Document
            doc = Document(str(p))
            text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        elif ext == ".doc":
            text = _read_doc_with_com(p)
        else:
            log.warning(f"不支持的教案格式 {ext}：{path}")
            return ""
        text = text.strip()
        if len(text) > 8000:
            text = text[:8000]
            log.info(f"教案过长，已截断到 8000 字（原始 {len(text)} 字）")
        return text
    except Exception as e:
        log.warning(f"教案文本提取失败 {path}：{e}")
        return ""


# 教案标题正则：第X章/节、一、、1.、1、、# 标题、【标题】
_LESSON_HEADING_RE = re.compile(
    r"^(?:"
    r"第[一二三四五六七八九十百千]+[章节部分]"
    r"|[一二三四五六七八九十]+、"
    r"|\d+[\.、)]"
    r"|\#{1,6}\s+\S"
    r"|【[^】]+】"
    r")"
)


def extract_lesson_plan_sections(path: str) -> List[tuple]:
    """从教案文件提取分段文本，返回 [(section_title, section_text), ...]。

    按常见教案标题格式分段（第X章/节、一、、1.、# 标题、【标题】）。
    无标题的文本归入"正文"段。用于按段匹配 PPT 页，避免整篇发给 LLM。
    """
    p = Path(path)
    if not p.exists():
        log.warning(f"教案文件不存在：{path}")
        return []
    ext = p.suffix.lower()
    try:
        if ext in (".txt", ".md"):
            try:
                text = p.read_text(encoding="utf-8")
            except UnicodeDecodeError:
                text = p.read_text(encoding="gbk")
        elif ext == ".docx":
            from docx import Document
            doc = Document(str(p))
            text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
        elif ext == ".doc":
            text = _read_doc_with_com(p)
        else:
            log.warning(f"不支持的教案格式 {ext}：{path}")
            return []
    except Exception as e:
        log.warning(f"教案分段提取失败 {path}：{e}")
        return []

    sections: List[tuple] = []
    cur_title = "正文"
    cur_lines: List[str] = []
    for line in text.splitlines():
        line_stripped = line.strip()
        if not line_stripped:
            if cur_lines:
                cur_lines.append("")
            continue
        if _LESSON_HEADING_RE.match(line_stripped):
            if cur_lines:
                sections.append((cur_title, "\n".join(cur_lines).strip()))
            cur_title = line_stripped
            cur_lines = []
        else:
            cur_lines.append(line_stripped)
    if cur_lines:
        sections.append((cur_title, "\n".join(cur_lines).strip()))

    if not sections:
        # 无任何内容
        return []
    log.info(f"教案分段完成：{len(sections)} 段")
    return sections


def _tokenize_zh(text: str) -> set:
    """简易中文分词：按非字母数字字符分割，过滤单字符停用词。"""
    # 按非字母数字（含中文标点、空格）分割
    tokens = re.split(r"[^\w]+", text)
    return {t.lower() for t in tokens if len(t) >= 2}


def match_lesson_sections(scene: "Scene", sections: List[tuple]) -> str:
    """选出与当前 PPT 页最相关的 1-2 个教案段，拼接后截断到 2000 字。

    用 Jaccard 相似度（关键词集合交集/并集）匹配。全部为 0 时回退前 2000 字。
    """
    if not sections:
        return ""
    scene_text = scene.title + " " + " ".join(scene.bullets)
    scene_tokens = _tokenize_zh(scene_text)
    if not scene_tokens:
        # 场景无关键词，回退前 2000 字
        return "\n\n".join(f"【{t}】\n{c}" for t, c in sections)[:2000]

    scored = []
    for idx, (title, content) in enumerate(sections):
        section_tokens = _tokenize_zh(title + " " + content)
        if not section_tokens:
            continue
        inter = scene_tokens & section_tokens
        union = scene_tokens | section_tokens
        score = len(inter) / len(union) if union else 0.0
        scored.append((score, idx, title, content))

    scored.sort(key=lambda x: x[0], reverse=True)
    # 取相似度 > 0 的前 2 段
    top = [s for s in scored if s[0] > 0][:2]
    if not top:
        # 全部为 0，回退前 2000 字
        return "\n\n".join(f"【{t}】\n{c}" for t, c in sections)[:2000]

    parts = [f"【{title}】\n{content}" for _, _, title, content in top]
    result = "\n\n".join(parts)
    if len(result) > 2000:
        result = result[:2000]
    return result


@dataclass
class Scene:
    index: int
    title: str
    bullets: List[str] = field(default_factory=list)
    image_path: Optional[str] = None
    script: str = ""
    audio_path: Optional[str] = None
    duration_sec: float = 0.0


@dataclass
class Job:
    job_id: str
    filename: str
    pptx_path: str
    status: str = JOB_STATUS_PENDING
    progress: float = 0.0
    stage: str = "等待中"
    avatar: str = "teacher_female"
    voice: str = "zh-CN-XiaoxiaoNeural"
    ratio: str = "16:9"
    resolution: str = "720p"
    digital_human_mode: str = "auto"  # auto=sadtalker优先 | static=静态头像
    enable_subtitle: bool = True
    enable_bgm: bool = True
    lesson_plan_path: Optional[str] = None  # 教案文件路径（可选，传给 LLM 参考生成文案）
    scenes: List[Scene] = field(default_factory=list)
    output_path: Optional[str] = None
    error: Optional[str] = None
    created_at: float = field(default_factory=time.time)
    started_at: Optional[float] = None
    finished_at: Optional[float] = None

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        return d


def _cleanup_job_files(job: "Job") -> List[str]:
    """清理任务相关的所有文件和目录。

    从 job.pptx_path 和 job.output_path 推导所有产物路径并删除。
    单个路径失败不影响其他路径，返回已清理路径列表供日志记录。

    清理范围：
    - uploads/{stem}.pptx              原始上传文件
    - uploads/{stem}.pdf               PPT 转 PDF 产物
    - uploads/_mineru_{stem}/          MinerU 解析结果目录
    - uploads/_slides_{stem}/          PDF 转 PNG 切片目录
    - uploads/{教案文件}                教案上传文件（若存在）
    - output/{job_id}/                 输出工作目录（含 audio/ scenes/ final/ talking_head_*/）
    """
    import shutil
    cleaned: List[str] = []

    # 1) 从 pptx_path 推导 uploads/ 下所有产物
    if job.pptx_path:
        pptx = Path(job.pptx_path)
        stem = pptx.stem  # 如 d62a9c89_test
        uploads_dir = pptx.parent

        upload_artifacts = [
            pptx,                                  # 原始 PPTX
            uploads_dir / f"{stem}.pdf",           # PPT 转 PDF
            uploads_dir / f"_mineru_{stem}",       # MinerU 解析目录
            uploads_dir / f"_slides_{stem}",       # PNG 切片目录
        ]
        # 教案文件（若存在）
        if job.lesson_plan_path:
            upload_artifacts.append(Path(job.lesson_plan_path))
        for p in upload_artifacts:
            try:
                if p.is_dir():
                    shutil.rmtree(p, ignore_errors=False)
                    cleaned.append(str(p))
                elif p.exists():
                    p.unlink(missing_ok=False)
                    cleaned.append(str(p))
            except Exception as e:
                log.warning(f"清理失败 {p}: {e}")

    # 2) 从 output_path 推导 output/{job_id}/ 整个目录
    if job.output_path:
        out_file = Path(job.output_path)
        # output_path = output/{job_id}/final/{stem}.mp4
        # 工作目录 = output/{job_id}/
        job_work_dir = out_file.parent.parent
        if job_work_dir.exists() and job_work_dir.is_dir():
            # 安全检查：必须是 OUTPUT_DIR 的直接子目录，避免误删
            try:
                if job_work_dir.parent == OUTPUT_DIR:
                    shutil.rmtree(job_work_dir, ignore_errors=False)
                    cleaned.append(str(job_work_dir))
                else:
                    log.warning(
                        f"拒绝清理非 OUTPUT_DIR 子目录: {job_work_dir}"
                    )
            except Exception as e:
                log.warning(f"清理失败 {job_work_dir}: {e}")

    return cleaned


class JobStore:
    """简单的 JSON 文件持久化的任务存储"""
    def __init__(self, path: Path = JOBS_FILE):
        self.path = path
        self.lock = threading.Lock()
        self._cache: Dict[str, Job] = {}
        self._load()

    def _load(self):
        if self.path.exists():
            try:
                with self.path.open("r", encoding="utf-8") as f:
                    data = json.load(f)
                for j in data.get("jobs", []):
                    # 重建 Scene 对象
                    scenes = [Scene(**s) for s in j.pop("scenes", [])]
                    self._cache[j["job_id"]] = Job(scenes=scenes, **j)
            except Exception as e:
                log.warning(f"加载 jobs.json 失败：{e}")

    def _persist(self):
        with self.path.open("w", encoding="utf-8") as f:
            data = {"jobs": [j.to_dict() for j in self._cache.values()]}
            json.dump(data, f, ensure_ascii=False, indent=2)

    def add(self, job: Job):
        with self.lock:
            self._cache[job.job_id] = job
            self._persist()

    def update(self, job_id: str, **kwargs):
        with self.lock:
            job = self._cache.get(job_id)
            if not job:
                return None
            for k, v in kwargs.items():
                setattr(job, k, v)
            self._persist()
            return job

    def get(self, job_id: str) -> Optional[Job]:
        with self.lock:
            return self._cache.get(job_id)

    def all(self) -> List[Job]:
        with self.lock:
            return sorted(self._cache.values(), key=lambda j: j.created_at, reverse=True)

    def delete(self, job_id: str):
        with self.lock:
            self._cache.pop(job_id, None)
            self._persist()


store = JobStore()


# ============ 1. PPT 解析 ============
def parse_pptx(pptx_path: str) -> List[Scene]:
    """解析 PPT 每页幻灯片：python-pptx 取标题，MinerU 解析 PDF 取正文+公式"""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("缺少 python-pptx，请先运行 pip install python-pptx")

    prs = Presentation(pptx_path)
    pptx_path = Path(pptx_path)

    # 1) PPT → PDF（PowerPoint COM 或 LibreOffice）
    pdf_path = _convert_pptx_to_pdf(pptx_path)
    if not pdf_path:
        raise RuntimeError("PPT 转 PDF 失败，无法解析")

    # 2) MinerU 解析整个 PDF，按页返回 [(text, kind, y), ...]
    pages_ocr = _ocr_pdf_with_mineru(str(pdf_path))

    # 3) 仍需每页 PNG（用于视频片段渲染）
    slide_imgs_dir = pptx_path.parent / f"_slides_{pptx_path.stem}"
    slide_imgs_dir.mkdir(exist_ok=True)
    _pdf_to_images(pdf_path, slide_imgs_dir)
    img_files = sorted(slide_imgs_dir.glob("slide_*.png"))

    scenes: List[Scene] = []
    for idx, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    title = text
                    break
            if title:
                break

        img_path = str(img_files[idx]) if idx < len(img_files) else None
        bullets: List[str] = []

        # 取该页的 OCR 结果
        ocr_items = pages_ocr[idx] if idx < len(pages_ocr) else []
        seen_norm = set()
        if title:
            seen_norm.add(title.replace(" ", ""))
        for text, kind, _y in ocr_items:
            norm = text.replace(" ", "")
            if norm in seen_norm:
                continue
            if kind == "formula":
                bullets.append(f"@@FORMULA@@{text}")
            else:
                bullets.append(text)
            seen_norm.add(norm)

        if not title:
            title = f"第 {idx + 1} 页"

        scenes.append(Scene(index=idx, title=title, bullets=bullets, image_path=img_path))

    return scenes


def _ppt_com_to_pdf(pptx_path: Path, pdf_path: Path) -> bool:
    """PowerPoint COM 自动化：PPT → PDF"""
    try:
        import win32com.client
        import pythoncom
    except ImportError:
        log.info("pywin32 未安装，跳过 PowerPoint COM 方案")
        return False

    abs_pptx = str(pptx_path.resolve())
    abs_pdf = str(pdf_path.resolve())
    app = None
    pres = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(abs_pptx, ReadOnly=True, WithWindow=False)
        pres.SaveAs(abs_pdf, 32)
        pres.Close()
        pres = None
        app.Quit()
        app = None
        if pdf_path.exists() and pdf_path.stat().st_size > 0:
            log.info(f"PowerPoint COM: {pdf_path.stat().st_size} bytes")
            return True
        return False
    except Exception as e:
        log.warning(f"PowerPoint COM 转换失败：{e}")
        return False
    finally:
        try:
            if pres:
                pres.Close()
        except Exception:
            pass
        try:
            if app:
                app.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def _soffice_to_pdf(pptx_path: Path, pdf_path: Path, timeout: int = 180) -> bool:
    """LibreOffice 转 PDF"""
    import tempfile
    tmp_profile = tempfile.mkdtemp(prefix="lo_profile_")
    profile_url = Path(tmp_profile).as_uri()

    cmd = [
        "soffice",
        f"-env:UserInstallation={profile_url}",
        "--headless", "--convert-to", "pdf",
        "--outdir", str(pdf_path.parent), str(pptx_path)
    ]

    try:
        proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    except FileNotFoundError:
        log.warning("未找到 soffice (LibreOffice)")
        return False

    start = time.time()
    while True:
        elapsed = time.time() - start
        if proc.poll() is not None:
            break
        if pdf_path.exists() and pdf_path.stat().st_size > 0:
            time.sleep(1)
            if pdf_path.exists() and pdf_path.stat().st_size > 0:
                proc.terminate()
                try:
                    proc.wait(timeout=5)
                except Exception:
                    proc.kill()
                shutil.rmtree(tmp_profile, ignore_errors=True)
                log.info(f"LibreOffice: PDF 生成成功（{elapsed:.0f}s）")
                return True
        if elapsed > timeout:
            log.warning(f"LibreOffice 转换超时（{timeout}s）")
            proc.kill()
            try:
                proc.wait(timeout=5)
            except Exception:
                pass
            shutil.rmtree(tmp_profile, ignore_errors=True)
            if pdf_path.exists() and pdf_path.stat().st_size > 0:
                log.info("LibreOffice: 超时但 PDF 已生成")
                return True
            return False
        time.sleep(2)

    shutil.rmtree(tmp_profile, ignore_errors=True)
    if pdf_path.exists() and pdf_path.stat().st_size > 0:
        log.info(f"LibreOffice: PDF 生成成功（{time.time()-start:.0f}s）")
        return True
    log.warning(f"LibreOffice: 进程退出但无 PDF（exit={proc.returncode}）")
    return False


def _convert_pptx_to_pdf(pptx_path: Path) -> Optional[Path]:
    """PPT → PDF：PowerPoint COM → LibreOffice，返回 PDF 路径，失败返回 None"""
    pdf_path = pptx_path.parent / (pptx_path.stem + ".pdf")

    try:
        subprocess.run(["taskkill", "/F", "/IM", "soffice.exe"],
                       capture_output=True, timeout=5)
        subprocess.run(["taskkill", "/F", "/IM", "soffice.bin"],
                       capture_output=True, timeout=5)
    except Exception:
        pass

    com_ok = _ppt_com_to_pdf(pptx_path, pdf_path)
    if com_ok:
        return pdf_path

    soff_ok = _soffice_to_pdf(pptx_path, pdf_path, timeout=180)
    if soff_ok:
        time.sleep(1)
        return pdf_path

    return None


def _pdf_to_images(pdf_path: Path, out_dir: Path):
    """PDF 每页 → PNG"""
    try:
        import pypdfium2 as pdfium
        doc = pdfium.PdfDocument(str(pdf_path))
        for i in range(len(doc)):
            bitmap = doc[i].render(scale=2.0)
            bitmap.to_pil().save(out_dir / f"slide_{i:03d}.png", "PNG")
        log.info(f"pypdfium2: PDF → {len(doc)} 张 PNG")
        return
    except ImportError:
        log.warning("pypdfium2 未安装")
    except Exception as e:
        log.warning(f"pypdfium2 渲染失败：{e}")

    try:
        from pdf2image import convert_from_path
        pages = convert_from_path(str(pdf_path), dpi=150)
        for i, page in enumerate(pages):
            page.save(out_dir / f"slide_{i:03d}.png", "PNG")
        log.info(f"pdf2image: PDF → {len(pages)} 张 PNG")
    except Exception as e:
        log.warning(f"pdf 转图片失败：{e}")


# ============ 2. 讲解稿生成 ============
_SCRIPT_LLM_CACHE: Dict[str, str] = {}


def generate_script_with_lesson_plan(scene: Scene, sections: List[tuple]) -> str:
    """参考教案用 DeepSeek 把 PPT 要点扩展为完整口播文案。

    接收教案分段列表 [(title, text), ...]，内部按 Jaccard 相似度选出与当前 PPT 页
    最相关的 1-2 段喂给 LLM，避免整篇教案发送导致 token 浪费和相关性下降。
    失败时回退到 generate_script(scene) 保持原有行为。
    公式占位符 @@FORMULA@@ 在 LLM 阶段替换为 [公式: ...]，LLM 只生成文本部分，
    公式仍由 latex_to_chinese_llm 单独处理。
    """
    if not sections:
        return generate_script(scene)

    api_key = os.environ.get("DEEPSEEK_API_KEY", "").strip()
    if not api_key:
        return generate_script(scene)

    # 匹配与本页最相关的教案段
    matched_text = match_lesson_sections(scene, sections)
    if not matched_text:
        return generate_script(scene)

    # 缓存 key：场景索引 + 匹配文本 hash + bullets hash
    cache_key = f"{scene.index}_{hash(matched_text)}_{hash(tuple(scene.bullets))}"
    if cache_key in _SCRIPT_LLM_CACHE:
        return _SCRIPT_LLM_CACHE[cache_key]

    # bullets 预处理：公式项转占位符，让 LLM 知道是公式但不读
    bullets_for_llm = []
    for b in scene.bullets:
        if b.startswith("@@FORMULA@@"):
            latex = b[len("@@FORMULA@@"):]
            bullets_for_llm.append(f"[公式: {latex}]")
        else:
            bullets_for_llm.append(b)
    bullets_text = "\n".join(f"- {b}" for b in bullets_for_llm if b)

    system_prompt = (
        "你是教学文案撰写助手。我会给你一段与当前 PPT 页相关的教案内容和这页 PPT 的要点，"
        "请参考教案内容，把 PPT 要点扩展成自然流畅的口播文案。\n\n"
        "要求：\n"
        "1. 文案要口语化，适合朗读，避免书面语\n"
        "2. 可以参考教案补充背景知识、过渡衔接、举例说明\n"
        "3. 保持 PPT 原要点的核心信息，不要遗漏\n"
        "4. 公式占位符 [公式: ...] 原样保留，不要尝试读公式\n"
        "5. 单页文案控制在 300 字以内\n"
        "6. 不要加\"大家好\"\"接下来\"等与 PPT 无关的套话\n"
        "7. 直接输出文案，不要加解释说明"
    )
    user_content = (
        f"【相关教案段落】\n{matched_text}\n\n"
        f"【本页要点】\n标题：{scene.title}\n{bullets_text}"
    )

    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com/v1", timeout=60)
        resp = client.chat.completions.create(
            model="deepseek-chat",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_content},
            ],
            temperature=0.3,
            max_tokens=800,
        )
        raw = resp.choices[0].message.content.strip()
        if not raw:
            log.warning("DeepSeek 文案生成返回空，回退 generate_script")
            return generate_script(scene)

        # 把 [公式: ...] 占位符替换回 @@FORMULA@@，再走原有公式转换流程
        def _restore_formula(m):
            return f"@@FORMULA@@{m.group(1)}"

        raw = re.sub(r"\[公式:\s*(.*?)\s*\]", _restore_formula, raw)

        # 按原有流程处理：拆分 parts，公式走 latex_to_chinese_llm，文本走 math_symbols_to_chinese
        parts = []
        # 用 @@FORMULA@@ 作为分隔标记切分
        segments = re.split(r"(@@FORMULA@@[^\n]*)", raw)
        for seg in segments:
            if not seg.strip():
                continue
            if seg.startswith("@@FORMULA@@"):
                latex = seg[len("@@FORMULA@@"):]
                latex_fixed = latex.replace(r"\nu", "v")
                spoken = latex_to_chinese_llm(latex_fixed)
                spoken = math_symbols_to_chinese(spoken)
                if spoken:
                    parts.append(spoken + "。")
            else:
                # 普通文本段，按句号拆分后逐句清洗
                for sentence in re.split(r"[。！？\n]", seg):
                    s = sentence.strip().rstrip(".。")
                    if s:
                        parts.append(math_symbols_to_chinese(s) + "。")

        text = " ".join(parts) if parts else raw
        if len(text) > 400:
            text = text[:400] + "..."
        _SCRIPT_LLM_CACHE[cache_key] = text
        log.info(f"[scene {scene.index}] 教案 LLM 文案生成成功（{len(text)} 字）")
        return text
    except Exception as e:
        log.warning(f"DeepSeek 教案文案生成失败，回退 generate_script：{e}")
        return generate_script(scene)


def generate_script(scene: Scene) -> str:
    """把 title + bullets 拼成口播文案，公式转中文口播，超长截断不跨页"""
    parts = [math_symbols_to_chinese(scene.title) + "。"]
    for b in scene.bullets:
        if not b:
            continue
        if b.startswith("@@FORMULA@@"):
            latex = b[len("@@FORMULA@@"):]
            latex_fixed = latex.replace(r"\nu", "v")
            spoken = latex_to_chinese_llm(latex_fixed)
            spoken = math_symbols_to_chinese(spoken)
            if spoken:
                parts.append(spoken + "。")
        else:
            b = b.strip().rstrip("。.")
            if b:
                parts.append(math_symbols_to_chinese(b) + "。")
    text = " ".join(parts)
    if len(text) > 400:
        text = text[:400] + "..."
    return text


# ============ 3. TTS 合成 ============
def tts_synthesize(text: str, voice: str, out_path: str) -> float:
    """edge-tts 合成语音，返回时长（秒），失败回退 pyttsx3"""
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        import edge_tts
        import asyncio

        async def _run():
            communicate = edge_tts.Communicate(text, voice)
            await communicate.save(str(out_path))

        asyncio.run(_run())
        if out_path.exists() and out_path.stat().st_size > 0:
            return _audio_duration(out_path)
    except Exception as e:
        log.warning(f"edge-tts 失败：{e}")

    try:
        import pyttsx3
        engine = pyttsx3.init()
        engine.save_to_file(text, str(out_path))
        engine.runAndWait()
        if out_path.exists() and out_path.stat().st_size > 0:
            return _audio_duration(out_path)
    except Exception as e:
        log.warning(f"pyttsx3 失败：{e}")

    log.warning("TTS 全部失败，生成静音占位")
    _make_silent_audio(out_path, duration=3.0)
    return 3.0


def _audio_duration(audio_path: Path) -> float:
    """用 ffmpeg 探测音频时长"""
    try:
        result = subprocess.run(
            [FFMPEG_BIN, "-i", str(audio_path), "-f", "null", "-"],
            capture_output=True, text=True, timeout=10
        )
        m = re.search(r"Duration:\s*(\d+):(\d+):(\d+(?:\.\d+)?)", result.stderr)
        if m:
            h, mi, s = int(m.group(1)), int(m.group(2)), float(m.group(3))
            return h * 3600 + mi * 60 + s
    except Exception:
        pass
    return 3.0


def _make_silent_audio(out_path: Path, duration: float):
    """生成静音 MP3"""
    try:
        subprocess.run([
            FFMPEG_BIN, "-y", "-f", "lavfi",
            "-i", f"anullsrc=channel_layout=stereo:sample_rate=44100",
            "-t", str(duration), "-q:a", "9", "-acodec", "libmp3lame",
            str(out_path)
        ], check=True, capture_output=True)
    except Exception as e:
        log.error(f"无法生成静音音频：{e}")


# ============ 4. 数字人视频合成 ============

def _is_sadtalker_available() -> bool:
    sadtalker_dir = BASE_DIR / "sadtalker"
    return (sadtalker_dir / "inference.py").exists()


def _generate_talking_head(avatar_img: str, audio_path: str,
                           result_dir: str) -> Optional[str]:
    """SadTalker 对口型数字人视频，失败返回 None"""
    if not _is_sadtalker_available():
        return None

    wrapper = BASE_DIR / "digital_human" / "sadtalker_wrapper.py"
    if not wrapper.exists():
        return None

    Path(result_dir).mkdir(parents=True, exist_ok=True)

    cmd = [
        sys.executable, str(wrapper),
        "--source_image", str(avatar_img),
        "--driven_audio", str(audio_path),
        "--result_dir", str(result_dir),
        "--full_body"
    ]

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=900)
    except subprocess.TimeoutExpired:
        log.error("[SadTalker] 生成超时")
        return None
    except Exception as e:
        log.error(f"[SadTalker] 调用异常: {e}")
        return None

    if result.returncode != 0:
        log.error(f"[SadTalker] 生成失败: {result.stderr[:500] if result.stderr else '未知'}")
        return None

    for line in result.stdout.splitlines():
        if line.startswith("RESULT:"):
            video_path = line[7:].strip()
            if Path(video_path).exists():
                return video_path

    mp4_files = sorted(
        Path(result_dir).rglob("*.mp4"),
        key=lambda f: f.stat().st_mtime, reverse=True
    )
    return str(mp4_files[0]) if mp4_files else None


def render_scene_video(scene: Scene, out_path: str, avatar: str,
                       ratio: str, resolution: str,
                       digital_human_mode: str = "auto") -> str:
    """合成 PPT 页面 + 数字人 + 语音 为一段视频"""
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    w, h = _ratio_size(ratio, resolution)

    avatar_img = AVATAR_DIR / f"{avatar}.png"
    if not avatar_img.exists():
        avatar_img = _make_placeholder_avatar(avatar, w, h)

    bg = scene.image_path if scene.image_path and Path(scene.image_path).exists() else None
    audio_dur = max(scene.duration_sec, 2.0)
    audio_path = scene.audio_path if scene.audio_path and Path(scene.audio_path).exists() else None

    talking_head_video = None
    if digital_human_mode == "auto" and audio_path and _is_sadtalker_available():
        talking_head_video = _generate_talking_head(
            str(avatar_img), audio_path,
            str(out_path.parent / f"talking_head_{scene.index}")
        )

    if talking_head_video:
        if bg:
            bg_input = ["-loop", "1", "-t", str(audio_dur), "-i", str(bg)]
        else:
            bg_input = [
                "-f", "lavfi", "-t", str(audio_dur),
                "-i", f"color=c=0x1a1f3a:s={w}x{h}:r=25"
            ]

        avatar_w = int(w * 0.28)
        avatar_h = int(avatar_w * 0.75)
        avatar_x = w - avatar_w - int(w * 0.03)
        avatar_y = h - avatar_h - int(h * 0.03)

        if audio_path:
            audio_input = ["-i", str(audio_path)]
            audio_map = "2:a"
        else:
            audio_input = [
                "-f", "lavfi", "-t", str(audio_dur),
                "-i", "anullsrc=channel_layout=stereo:sample_rate=44100"
            ]
            audio_map = "2:a"

        filter_complex = (
            f"[0:v]scale={w}:{h}:force_original_aspect_ratio=decrease,"
            f"pad={w}:{h}:(ow-iw)/2:(oh-ih)/2:black,setsar=1,fps=25[bg];"
            f"[1:v]scale={avatar_w}:{avatar_h},fps=25,format=yuva420p[avt];"
            f"[bg][avt]overlay={avatar_x}:{avatar_y}[v]"
        )

        cmd = [
            FFMPEG_BIN, "-y",
            *bg_input,
            "-i", str(talking_head_video),
            *audio_input,
            "-filter_complex", filter_complex,
            "-map", "[v]", "-map", audio_map,
            "-c:v", "libx264", "-preset", "ultrafast", "-crf", "26",
            "-c:a", "aac", "-b:a", "96k",
            "-pix_fmt", "yuv420p", "-shortest",
            "-t", str(audio_dur),
            str(out_path)
        ]
    else:
        if bg:
            bg_input = ["-loop", "1", "-t", str(audio_dur), "-i", str(bg)]
        else:
            bg_input = [
                "-f", "lavfi", "-t", str(audio_dur),
                "-i", f"color=c=0x1a1f3a:s={w}x{h}:r=25"
            ]

        main_w = int(w * 0.75)
        bg_h = int(main_w * h / w)
        bg_y = (h - bg_h) // 2
        side_w = w - main_w
        avatar_w = int(side_w * 0.7)
        avatar_h = avatar_w
        avatar_x = main_w + (side_w - avatar_w) // 2
        avatar_y = (h - avatar_h) // 2
        avatar_input = ["-loop", "1", "-t", str(audio_dur), "-i", str(avatar_img)]

        if audio_path:
            audio_input = ["-i", str(audio_path)]
            audio_map = "2:a"
        else:
            audio_input = [
                "-f", "lavfi", "-t", str(audio_dur),
                "-i", "anullsrc=channel_layout=stereo:sample_rate=44100"
            ]
            audio_map = "2:a"

        filter_complex = (
            f"[0:v]scale={main_w}:{bg_h}:force_original_aspect_ratio=decrease,"
            f"pad={main_w}:{bg_h}:(ow-iw)/2:(oh-ih)/2:black,setsar=1,fps=25[bg1];"
            f"[bg1]pad={w}:{h}:0:{bg_y}:0x1a1f3a[bg];"
            f"[1:v]scale={avatar_w}:{avatar_h},fps=25[avt];"
            f"[bg][avt]overlay={avatar_x}:{avatar_y}[v]"
        )

        cmd = [
            FFMPEG_BIN, "-y",
            *bg_input,
            *avatar_input,
            *audio_input,
            "-filter_complex", filter_complex,
            "-map", "[v]", "-map", audio_map,
            "-c:v", "libx264", "-preset", "ultrafast", "-crf", "26",
            "-c:a", "aac", "-b:a", "96k",
            "-pix_fmt", "yuv420p", "-shortest",
            "-t", str(audio_dur),
            str(out_path)
        ]

    try:
        subprocess.run(cmd, check=True, capture_output=True, timeout=180)
    except subprocess.CalledProcessError as e:
        err = e.stderr.decode(errors='ignore') if e.stderr else ''
        log.error(f"FFmpeg 渲染失败 (exit {e.returncode})：{err[:1000]}")
        raise
    except FileNotFoundError:
        raise RuntimeError("未找到 ffmpeg")

    return str(out_path)


def _ratio_size(ratio: str, resolution: str) -> tuple:
    presets = {
        ("16:9", "720p"): (1280, 720),
        ("16:9", "1080p"): (1920, 1080),
        ("9:16", "720p"): (720, 1280),
        ("9:16", "1080p"): (1080, 1920),
        ("1:1", "720p"): (720, 720),
        ("1:1", "1080p"): (1080, 1080),
    }
    return presets.get((ratio, resolution), (1280, 720))


def _make_placeholder_avatar(name: str, w: int, h: int) -> Path:
    """生成占位头像"""
    png_path = AVATAR_DIR / f"{name}.png"
    if png_path.exists():
        return png_path
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return png_path

    palettes = {
        "teacher_female": ("#6c8cff", "#3a4ba0", "教师·女"),
        "teacher_male":   ("#3a4ba0", "#1f2a6b", "教师·男"),
        "young_girl":     ("#ff6ba0", "#b0306b", "少女"),
        "young_boy":      ("#36d6c0", "#1f8a7a", "少年"),
        "professor":      ("#b06bff", "#5a2da0", "教授"),
    }
    c1, c2, label = palettes.get(name, ("#6c8cff", "#3a4ba0", name))

    img = Image.new("RGB", (400, 400), c1)
    d = ImageDraw.Draw(img)
    r1, g1, b1 = int(c1[1:3], 16), int(c1[3:5], 16), int(c1[5:7], 16)
    r2, g2, b2 = int(c2[1:3], 16), int(c2[3:5], 16), int(c2[5:7], 16)
    for y in range(400):
        t = y / 400
        r = int(r1 + (r2 - r1) * t)
        g = int(g1 + (g2 - g1) * t)
        b = int(b1 + (b2 - b1) * t)
        d.line([(0, y), (400, y)], fill=(r, g, b))
    d.ellipse([130, 100, 270, 240], fill="#ffe2c8")
    d.polygon([(100, 280), (300, 280), (320, 400), (80, 400)], fill="#ffe2c8")
    d.ellipse([165, 165, 180, 180], fill="#1a1a1a")
    d.ellipse([220, 165, 235, 180], fill="#1a1a1a")
    d.arc([180, 195, 220, 220], start=0, end=180, fill="#1a1a1a", width=4)
    try:
        font = ImageFont.truetype("arial.ttf", 32)
    except Exception:
        font = ImageFont.load_default()
    d.text((200, 30), label, fill="white", anchor="mt", font=font)
    img.save(png_path, "PNG")
    return png_path


# ============ 5. 合并分镜 ============
def concatenate_scenes(scene_videos: List[str], audio_track: Optional[str],
                       bgm: Optional[str], out_path: str,
                       subtitle: bool = False) -> str:
    """拼接分镜视频，可选叠加 BGM"""
    if not scene_videos:
        raise ValueError("没有可拼接的分镜视频")

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    concat_list = out_path.parent / "_concat.txt"
    with concat_list.open("w", encoding="utf-8") as f:
        for v in scene_videos:
            f.write(f"file '{v.replace(chr(92), '/')}'\n")

    tmp_concat = out_path.with_suffix(".concat.mp4")
    subprocess.run([
        FFMPEG_BIN, "-y", "-f", "concat", "-safe", "0",
        "-i", str(concat_list), "-c", "copy", str(tmp_concat)
    ], check=True, capture_output=True)

    if bgm and Path(bgm).exists():
        subprocess.run([
            FFMPEG_BIN, "-y",
            "-i", str(tmp_concat), "-i", str(bgm),
            "-filter_complex",
            f"[0:a]volume=0.9[a0];[1:a]volume=0.15,afade=t=in:st=0:d=2,afade=t=out:st=-3:d=3[a1];"
            f"[a0][a1]amix=inputs=2:duration=first[aout]",
            "-map", "0:v", "-map", "[aout]",
            "-c:v", "copy", "-c:a", "aac", "-b:a", "192k",
            "-shortest", str(out_path)
        ], check=True, capture_output=True)
        tmp_concat.unlink(missing_ok=True)
    else:
        tmp_concat.rename(out_path)

    return str(out_path)


# ============ 6. 任务主流程 ============
def run_job(job_id: str, executor: ThreadPoolExecutor):
    """任务入口（后台线程池执行）"""
    job = store.get(job_id)
    if not job:
        return
    job.started_at = time.time()
    store.update(job_id, status=JOB_STATUS_RUNNING, stage="解析 PPT")

    # 每次任务创建调试文档目录，记录 PPT 提取 / 教案提取 / LLM 文案合成结果
    debug_job_dir = DEBUG_DIR / f"{job_id}_{time.strftime('%Y%m%d_%H%M%S')}"
    debug_job_dir.mkdir(parents=True, exist_ok=True)

    try:
        log.info(f"[{job_id}] 开始解析 PPT")
        scenes = parse_pptx(job.pptx_path)
        job.scenes = scenes
        store.update(job_id, scenes=scenes, progress=0.15, stage=f"已解析 {len(scenes)} 页分镜")

        # 记录 PPT 提取结果
        ppt_lines = [
            f"PPT 文件：{job.filename}",
            f"总页数：{len(scenes)}",
            f"提取时间：{time.strftime('%Y-%m-%d %H:%M:%S')}",
            "=" * 60,
            "",
        ]
        for sc in scenes:
            ppt_lines.append(f"【第 {sc.index + 1} 页】标题：{sc.title}")
            for b in sc.bullets:
                if b.startswith("@@FORMULA@@"):
                    ppt_lines.append(f"  - [公式] {b[len('@@FORMULA@@'):]}")
                else:
                    ppt_lines.append(f"  - {b}")
            ppt_lines.append("")
        (debug_job_dir / "01_ppt_extract.txt").write_text(
            "\n".join(ppt_lines), encoding="utf-8")

        log.info(f"[{job_id}] 生成讲解稿并合成 TTS")
        work_dir = OUTPUT_DIR / job_id
        audio_dir = work_dir / "audio"
        video_dir = work_dir / "scenes"
        final_dir = work_dir / "final"
        for d in [work_dir, audio_dir, video_dir, final_dir]:
            d.mkdir(parents=True, exist_ok=True)

        # 提取教案分段（若有），供 LLM 按段匹配生成文案
        lesson_sections: List[tuple] = []
        if job.lesson_plan_path and Path(job.lesson_plan_path).exists():
            lesson_sections = extract_lesson_plan_sections(job.lesson_plan_path)
            if lesson_sections:
                log.info(f"[{job_id}] 教案已加载：{len(lesson_sections)} 段")
            else:
                log.warning(f"[{job_id}] 教案分段提取为空，将走默认文案流程")

        # 记录教案提取结果
        lp_lines = [
            f"教案文件：{Path(job.lesson_plan_path).name if job.lesson_plan_path else '（未上传）'}",
            f"提取时间：{time.strftime('%Y-%m-%d %H:%M:%S')}",
            "=" * 60,
            "",
        ]
        if lesson_sections:
            lp_lines.append(f"共分段 {len(lesson_sections)} 段：")
            lp_lines.append("")
            for idx, (title, text) in enumerate(lesson_sections):
                lp_lines.append(f"【段 {idx + 1}】{title}（{len(text)} 字）")
                lp_lines.append(text)
                lp_lines.append("")
        else:
            lp_lines.append("未提取到教案分段（将走默认文案流程，不调用 LLM）")
        (debug_job_dir / "02_lesson_plan.txt").write_text(
            "\n".join(lp_lines), encoding="utf-8")

        # 记录每页最终文案
        script_lines = [
            f"任务 ID：{job_id}",
            f"PPT 文件：{job.filename}",
            f"文案路径：{'教案 LLM 生成' if lesson_sections else '默认规则生成'}",
            f"生成时间：{time.strftime('%Y-%m-%d %H:%M:%S')}",
            "=" * 60,
            "",
        ]

        for i, scene in enumerate(scenes):
            if store.get(job_id).status == JOB_STATUS_CANCELED:
                raise RuntimeError("任务被取消")
            used_llm = bool(lesson_sections)
            if used_llm:
                scene.script = generate_script_with_lesson_plan(scene, lesson_sections)
            else:
                scene.script = generate_script(scene)
            # 追加本页文案记录
            script_lines.append(f"【第 {scene.index + 1} 页】{scene.title}")
            script_lines.append(f"生成路径：{'教案 LLM' if used_llm else '默认规则'}")
            script_lines.append(f"字数：{len(scene.script)}")
            script_lines.append("文案内容：")
            script_lines.append(scene.script)
            script_lines.append("-" * 60)
            script_lines.append("")
            # 每页生成后立即写入，避免中途失败丢失记录
            (debug_job_dir / "03_llm_scripts.txt").write_text(
                "\n".join(script_lines), encoding="utf-8")

            audio_file = audio_dir / f"scene_{i:03d}.mp3"
            scene.duration_sec = tts_synthesize(scene.script, job.voice, str(audio_file))
            scene.audio_path = str(audio_file)
            store.update(job_id, scenes=scenes,
                         progress=0.15 + 0.25 * (i + 1) / len(scenes),
                         stage=f"TTS 第 {i + 1}/{len(scenes)} 页")

        log.info(f"[{job_id}] 渲染分镜视频")
        scene_videos: List[str] = []
        for i, scene in enumerate(scenes):
            if store.get(job_id).status == JOB_STATUS_CANCELED:
                raise RuntimeError("任务被取消")
            v = video_dir / f"scene_{i:03d}.mp4"
            render_scene_video(scene, str(v), job.avatar, job.ratio,
                               job.resolution, job.digital_human_mode)
            scene_videos.append(str(v))
            store.update(job_id,
                         progress=0.4 + 0.45 * (i + 1) / len(scenes),
                         stage=f"渲染第 {i + 1}/{len(scenes)} 页")

        log.info(f"[{job_id}] 合并成片")
        bgm_file = BGM_DIR / "default.mp3" if job.enable_bgm else None
        if bgm_file and not bgm_file.exists():
            bgm_file = None
        final = final_dir / f"{Path(job.filename).stem}.mp4"
        concatenate_scenes(scene_videos, None, str(bgm_file) if bgm_file else None,
                            str(final), subtitle=job.enable_subtitle)

        job.output_path = str(final)
        job.finished_at = time.time()
        store.update(job_id, status=JOB_STATUS_DONE, output_path=job.output_path,
                     progress=1.0, stage="完成", finished_at=job.finished_at)
        log.info(f"[{job_id}] 完成 → {final}")

    except Exception as e:
        log.exception(f"[{job_id}] 失败：{e}")
        store.update(job_id, status=JOB_STATUS_FAILED, error=str(e), stage="失败")


def submit_job(job: Job, executor: ThreadPoolExecutor) -> str:
    store.add(job)
    executor.submit(run_job, job.job_id, executor)
    return job.job_id


# ============ 测试入口 ============
if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("用法：python core.py <pptx 文件>")
        sys.exit(1)
    pptx = sys.argv[1]
    ex = ThreadPoolExecutor(max_workers=2)
    job = Job(
        job_id=uuid.uuid4().hex[:8],
        filename=Path(pptx).name,
        pptx_path=str(Path(pptx).resolve()),
    )
    submit_job(job, ex)
    print(f"任务已提交：{job.job_id}")
    while True:
        j = store.get(job.job_id)
        print(f"  {j.status} | {j.progress:.0%} | {j.stage}")
        if j.status in (JOB_STATUS_DONE, JOB_STATUS_FAILED):
            break
        time.sleep(2)
